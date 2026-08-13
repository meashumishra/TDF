"""DOCX / XLSX / PPTX readers."""

from __future__ import annotations

import re
from pathlib import Path

from ..ir import Doc, Figure, Heading, KV, ListBlock, Para, Table

_STYLE_H = re.compile(r"^Heading\s*(\d)", re.I)


def read_docx(path: str | Path) -> Doc:
    import docx
    from docx.table import Table as DTable
    from docx.text.paragraph import Paragraph

    d = docx.Document(str(path))
    doc = Doc(source=str(path))
    if d.core_properties.title:
        doc.title = d.core_properties.title

    items: list[str] = []
    ordered = False

    def flush():
        nonlocal items, ordered
        if items:
            doc.add(ListBlock(items, ordered))
            items, ordered = [], False

    # Walk the body in document order so tables stay where they belong.
    body = d.element.body
    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            p = Paragraph(child, d)
            text = p.text.strip()
            style = (p.style.name or "") if p.style is not None else ""
            if not text:
                continue
            if m := _STYLE_H.match(style):
                flush()
                doc.add(Heading(int(m.group(1)), text))
            elif style.lower().startswith("title"):
                flush()
                if not doc.title:
                    doc.title = text
                else:
                    doc.add(Heading(1, text))
            elif "list" in style.lower():
                if not items:
                    ordered = "number" in style.lower()
                items.append(text)
            else:
                flush()
                doc.add(Para(text))
        elif tag == "tbl":
            flush()
            t = DTable(child, d)
            grid = [[c.text.strip().replace("\n", " ") for c in row.cells] for row in t.rows]
            grid = [g for g in grid if any(g)]
            if grid:
                doc.add(Table(grid[0], grid[1:]))
    flush()
    return doc


def read_xlsx(path: str | Path) -> Doc:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    doc = Doc(source=str(path), title=Path(path).stem)

    for ws in wb.worksheets:
        rows = []
        for r in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v).strip() for v in r]
            if any(vals):
                rows.append(vals)
        if not rows:
            continue
        # Trim all-empty trailing columns.
        width = max(len(r) for r in rows)
        rows = [r + [""] * (width - len(r)) for r in rows]
        keep = [c for c in range(width) if any(r[c] for r in rows)]
        rows = [[r[c] for c in keep] for r in rows]

        header, data = rows[0], rows[1:]
        if not data:
            header, data = [f"c{i+1}" for i in range(len(rows[0]))], rows
        doc.add(Table(header, data, caption=ws.title, group=ws.title))
    wb.close()
    return doc


def read_pptx(path: str | Path) -> Doc:
    from pptx import Presentation

    prs = Presentation(str(path))
    doc = Doc(source=str(path), title=Path(path).stem)

    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        doc.add(Heading(2, title or f"Slide {idx}"))

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                grid = [[c.text.strip().replace("\n", " ") for c in row.cells]
                        for row in shape.table.rows]
                grid = [g for g in grid if any(g)]
                if grid:
                    doc.add(Table(grid[0], grid[1:]))
                continue
            if getattr(shape, "has_chart", False):
                doc.add(Figure(f"chart: {shape.chart.chart_type}", kind="chart"))
                continue
            if shape.has_text_frame:
                bullets = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                if len(bullets) > 1:
                    doc.add(ListBlock(bullets))
                elif bullets:
                    doc.add(Para(bullets[0]))

        if slide.has_notes_slide and (nt := slide.notes_slide.notes_text_frame.text.strip()):
            doc.add(KV([("notes", " ".join(nt.split()))]))
    return doc
