"""Generate realistic sample documents covering the cases where Markdown hurts.

Deliberately spans the four pathologies the research surfaced:
  1. wide sorted tables (spreadsheet exports)
  2. repeated legal/policy boilerplate
  3. running page headers/footers in PDFs
  4. slide decks with repeated chrome
"""

from __future__ import annotations

import random
from pathlib import Path

OUT = Path(__file__).resolve().parent.parent / "samples"
random.seed(7)

REGIONS = ["North America", "EMEA", "APAC", "LATAM"]
PRODUCTS = ["Widget Assembly", "Gadget Pro", "Sprocket Mini", "Flange Deluxe",
            "Bearing Standard", "Coupler XL"]
STATUS = ["shipped", "pending", "cancelled", "backordered"]

CLAUSE_A = ("Notwithstanding any other provision of this Agreement, the Company shall not be "
            "liable for any indirect, incidental, special, consequential or punitive damages")
CLAUSE_B = ("subject to the terms and conditions set forth in Section 12.4 of this Agreement "
            "and any applicable Statement of Work")
CLAUSE_C = ("Each party represents and warrants that it has full corporate power and authority "
            "to enter into this Agreement")


def sales_rows(n=400):
    rows = []
    for i in range(n):
        rows.append([
            str(10_000 + i),
            random.choice(PRODUCTS),
            random.choice(REGIONS),
            f"2025-{random.randint(1,12):02d}-{random.randint(1,28):02d}",
            f"${random.randint(50, 90000):,}.00",
            str(random.randint(1, 40)),
            random.choice(STATUS),
            "USD",
        ])
    rows.sort(key=lambda r: (r[2], r[1], r[3]))
    return rows


HDR = ["order_id", "product", "region", "order_date", "amount", "qty", "status", "currency"]


def make_xlsx():
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Q1 Orders"
    ws.append(HDR)
    for r in sales_rows(300):
        ws.append(r)
    ws2 = wb.create_sheet("Q2 Orders")
    ws2.append(HDR)
    for r in sales_rows(300):
        ws2.append(r)
    ws3 = wb.create_sheet("Summary")
    ws3.append(["region", "orders", "revenue"])
    for reg in REGIONS:
        ws3.append([reg, random.randint(50, 200), f"${random.randint(100000, 900000):,}.00"])
    p = OUT / "sales_report.xlsx"
    wb.save(p)
    return p


def make_docx():
    import docx

    d = docx.Document()
    d.add_heading("Master Services Agreement", 0)
    d.add_paragraph("This Agreement is entered into as of 1 January 2026.")
    for n in range(1, 13):
        d.add_heading(f"{n}. Clause {n}", level=1)
        d.add_paragraph(f"{CLAUSE_A}, {CLAUSE_B}.")
        d.add_paragraph(f"{CLAUSE_C}. The provisions of this Section {n} shall survive "
                        f"termination of this Agreement.")
        d.add_paragraph(f"For the purposes of this Section {n}, {CLAUSE_B} shall be "
                        f"interpreted in accordance with the laws of the State of Delaware.")
        if n % 4 == 0:
            d.add_heading(f"{n}.1 Fee schedule", level=2)
            t = d.add_table(rows=1, cols=4)
            for i, h in enumerate(["service", "tier", "rate", "currency"]):
                t.rows[0].cells[i].text = h
            for tier in ["Bronze", "Silver", "Gold", "Platinum"]:
                cells = t.add_row().cells
                for i, v in enumerate(["Managed hosting", tier,
                                       f"${random.randint(1000, 9000):,}.00", "USD"]):
                    cells[i].text = v
    p = OUT / "services_agreement.docx"
    d.save(p)
    return p


def make_pdf():
    import pymupdf as fitz

    doc = fitz.open()
    body = ("The quarterly operating review covers revenue attainment, pipeline coverage and "
            "regional performance. Management believes the results reflect continued execution "
            "against the operating plan established at the start of the fiscal year. ")
    for page_no in range(1, 15):
        page = doc.new_page()
        page.insert_text((60, 40), "ACME CORPORATION - CONFIDENTIAL - INTERNAL USE ONLY",
                         fontsize=8)
        page.insert_text((60, 80), f"Section {page_no}: Operating Review", fontsize=17)
        y = 115
        for para in range(5):
            text = body + f"Paragraph {para + 1} of section {page_no}."
            rect = fitz.Rect(60, y, 540, y + 90)
            page.insert_textbox(rect, text, fontsize=10)
            y += 95
        page.insert_text((60, 760),
                         "ACME CORPORATION - CONFIDENTIAL - INTERNAL USE ONLY", fontsize=8)
        page.insert_text((480, 760), f"Page {page_no} of 14", fontsize=8)
    p = OUT / "operating_review.pdf"
    doc.save(p)
    doc.close()
    return p


def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    layout = prs.slide_layouts[1]
    for i in range(1, 15):
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = f"Fiscal 2026 Operating Review - Section {i}"
        tf = s.placeholders[1].text_frame
        tf.text = "ACME Corporation - Confidential - Do not distribute"
        for b in range(4):
            p = tf.add_paragraph()
            p.text = (f"Key takeaway {b + 1}: regional performance tracked ahead of plan "
                      f"driven by improved pipeline conversion in section {i}.")
            p.font.size = Pt(14)
    p = OUT / "quarterly_deck.pptx"
    prs.save(p)
    return p


def make_csv():
    import csv

    p = OUT / "orders.csv"
    with open(p, "w", newline="", encoding="utf-8") as fh:
        w = csv.writer(fh)
        w.writerow(HDR)
        w.writerows(sales_rows(500))
    return p


def make_html():
    rows = sales_rows(120)
    trs = "".join(
        "<tr>" + "".join(f"<td>{c}</td>" for c in r) + "</tr>" for r in rows
    )
    ths = "".join(f"<th>{h}</th>" for h in HDR)
    paras = "".join(
        f"<h2>Section {i}</h2><p>{CLAUSE_A}, {CLAUSE_B}.</p><p>{CLAUSE_C}.</p>"
        for i in range(1, 9)
    )
    html = (f"<html><head><title>Operations Handbook</title></head><body>"
            f"<h1>Operations Handbook</h1>{paras}"
            f"<h2>Order detail</h2><table><tr>{ths}</tr>{trs}</table></body></html>")
    p = OUT / "handbook.html"
    p.write_text(html, encoding="utf-8")
    return p


def make_md():
    parts = ["# Platform Runbook", ""]
    for i in range(1, 10):
        parts += [f"## Service {i}", "",
                  f"{CLAUSE_C}. This runbook section describes the on-call procedure for "
                  f"service {i}. {CLAUSE_B}.", "",
                  "| check | threshold | owner | severity |",
                  "|---|---|---|---|"]
        for c in ["latency p99", "error rate", "saturation", "queue depth", "cache hit"]:
            parts.append(f"| {c} | {random.randint(1, 99)}% | platform-team | "
                         f"{random.choice(['sev1','sev2','sev3'])} |")
        parts.append("")
    p = OUT / "runbook.md"
    p.write_text("\n".join(parts), encoding="utf-8")
    return p


if __name__ == "__main__":
    OUT.mkdir(exist_ok=True)
    for fn in (make_xlsx, make_docx, make_pdf, make_pptx, make_csv, make_html, make_md):
        print("wrote", fn())
